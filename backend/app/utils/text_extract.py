"""Extracts raw text from the file types the platform accepts."""
import os
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation


def extract_text(file_path: str, file_type: str) -> str:
    file_type = file_type.lower()
    if file_type == "pdf":
        return _extract_pdf(file_path)
    if file_type == "docx":
        return _extract_docx(file_path)
    if file_type == "pptx":
        return _extract_pptx(file_path)
    if file_type == "txt" or file_type == "md":
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    raise ValueError(f"Unsupported file type: {file_type}")


def _extract_pdf(path: str) -> str:
    text_parts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text() or ""
            text_parts.append(page_text)
    return "\n".join(text_parts)


def _extract_docx(path: str) -> str:
    doc = DocxDocument(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(path: str) -> str:
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        parts.append(line)
    return "\n".join(parts)


def detect_file_type(filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower().replace(".", "")
    return ext
